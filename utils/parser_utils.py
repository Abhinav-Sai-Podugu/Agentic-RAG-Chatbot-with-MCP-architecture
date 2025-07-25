import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import csv
import re


def clean_text(text):
    """Clean and normalize text"""
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text)
    # Remove empty lines
    text = text.strip()
    return text


def chunk_text(text, max_chunk_size=500, overlap=50):
    """Split text into meaningful chunks with overlap"""
    if not text or len(text) < max_chunk_size:
        return [text] if text.strip() else []

    # Split by sentences first
    sentences = re.split(r'[.!?]+', text)

    chunks = []
    current_chunk = ""

    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue

        # If adding this sentence would exceed max size, save current chunk
        if len(current_chunk) + len(sentence) > max_chunk_size and current_chunk:
            chunks.append(current_chunk.strip())

            # Start new chunk with overlap (last few words)
            words = current_chunk.split()
            overlap_words = words[-overlap // 10:] if len(words) > overlap // 10 else []
            current_chunk = " ".join(overlap_words) + " " + sentence
        else:
            current_chunk += " " + sentence if current_chunk else sentence

    # Add the last chunk
    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    return chunks


def parse_document(file_path):
    """Parse document and return meaningful chunks"""
    try:
        if file_path.endswith(".pdf"):
            reader = PdfReader(file_path)
            full_text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    full_text += page_text + "\n"

            # Clean and chunk the full text
            clean_full_text = clean_text(full_text)
            return chunk_text(clean_full_text)

        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            full_text = ""
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text += para.text + "\n"

            clean_full_text = clean_text(full_text)
            return chunk_text(clean_full_text)

        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            chunks = []

            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + " "

                if slide_text.strip():
                    clean_slide_text = clean_text(slide_text)
                    # For presentations, keep slide-based chunks but clean them
                    chunks.append(f"[Slide {slide_num + 1}] {clean_slide_text}")

            return chunks

        elif file_path.endswith(".csv"):
            chunks = []
            with open(file_path, newline='', encoding='utf-8') as csvfile:
                reader = csv.reader(csvfile)
                headers = next(reader, None)

                if headers:
                    # Create a summary chunk with headers
                    chunks.append(f"CSV Headers: {', '.join(headers)}")

                    # Group rows into meaningful chunks
                    row_batch = []
                    batch_size = 10  # Process 10 rows at a time

                    for row_num, row in enumerate(reader):
                        if row and any(cell.strip() for cell in row):  # Skip empty rows
                            row_text = " | ".join(f"{headers[i] if i < len(headers) else f'Col{i}'}: {cell}"
                                                  for i, cell in enumerate(row) if cell.strip())
                            row_batch.append(row_text)

                            if len(row_batch) >= batch_size:
                                chunks.append("\n".join(row_batch))
                                row_batch = []

                    # Add remaining rows
                    if row_batch:
                        chunks.append("\n".join(row_batch))

            return chunks

        elif file_path.endswith((".txt", ".md")):
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            clean_content = clean_text(content)

            # For markdown, try to split by headers first
            if file_path.endswith(".md"):
                sections = re.split(r'\n#+\s', content)
                chunks = []
                for section in sections:
                    section = section.strip()
                    if section:
                        # Further chunk large sections
                        if len(section) > 1000:
                            chunks.extend(chunk_text(section))
                        else:
                            chunks.append(section)
                return chunks
            else:
                # For plain text, use paragraph-based chunking
                paragraphs = content.split('\n\n')
                chunks = []
                current_chunk = ""

                for para in paragraphs:
                    para = clean_text(para)
                    if not para:
                        continue

                    if len(current_chunk) + len(para) > 800:
                        if current_chunk:
                            chunks.append(current_chunk)
                        current_chunk = para
                    else:
                        current_chunk += "\n\n" + para if current_chunk else para

                if current_chunk:
                    chunks.append(current_chunk)

                return chunks

        return []

    except Exception as e:
        print(f"Error parsing {file_path}: {str(e)}")
        return [f"Error parsing file {os.path.basename(file_path)}: {str(e)}"]